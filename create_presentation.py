from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_slide_with_gradient(prs, title, subtitle=None, layout_type='title'):
    """Add a slide with custom formatting"""
    if layout_type == 'title':
        slide_layout = prs.slide_layouts[6]  # Blank layout
    else:
        slide_layout = prs.slide_layouts[6]  # Blank layout
    
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        left = Inches(0.5)
        top = Inches(2.8)
        width = Inches(9)
        height = Inches(1)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, items, bg_color=(102, 126, 234)):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add content box
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(102, 126, 234)

# Title
left = Inches(0.5)
top = Inches(2.5)
width = Inches(9)
height = Inches(1.5)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "🛡️ CyberGuardian AI"
p = title_frame.paragraphs[0]
p.font.size = Pt(66)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Subtitle
left = Inches(0.5)
top = Inches(4)
width = Inches(9)
height = Inches(1)
subtitle_box = slide.shapes.add_textbox(left, top, width, height)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Your All-in-One AI Cybersecurity Companion"
p = subtitle_frame.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Tech stack
left = Inches(2)
top = Inches(5.5)
width = Inches(6)
height = Inches(1)
tech_box = slide.shapes.add_textbox(left, top, width, height)
tech_frame = tech_box.text_frame
tech_frame.text = "Next.js • Gemini AI • VirusTotal API • Tesseract.js"
p = tech_frame.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Slide 2: The Problem
items = [
    "✓ Cybersecurity threats are increasing exponentially",
    "✓ Small businesses and individuals lack access to enterprise-grade security tools",
    "✓ Phishing attacks cost organizations $4.9 billion annually",
    "✓ Average person cannot identify sophisticated cyber threats",
    "✓ Incident response requires expensive specialized knowledge",
    "✓ Security awareness training is boring and ineffective"
]
add_content_slide(prs, "🚨 The Problem", items, (240, 147, 92))

# Slide 3: Market Impact
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(79, 172, 254)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "📊 The Impact"
p = title_frame.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Stats
stats = [
    ("91%", "Of cyberattacks start with phishing"),
    ("$6T", "Global cybercrime damage by 2025"),
    ("93%", "Of breaches are preventable")
]

for i, (number, text) in enumerate(stats):
    left = Inches(0.5 + i * 3.2)
    top = Inches(2.5)
    width = Inches(2.8)
    height = Inches(3)
    
    stat_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = stat_box.text_frame
    
    # Number
    p = text_frame.paragraphs[0]
    p.text = number
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)

# Slide 4: Our Solution
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(67, 233, 123)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "💡 Our Solution"
p = title_frame.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "AI-Powered, Accessible, All-in-One Platform"
p = subtitle_frame.paragraphs[0]
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Features
features = [
    ("📚 Cyber Awareness", "Interactive gamified learning"),
    ("🎣 Phishing Detection", "Real-time AI URL analysis"),
    ("🚨 Incident Response", "Step-by-step playbooks"),
    ("💻 Code Security", "Vulnerability scanning")
]

for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 4.7)
    top = Inches(2.5 + row * 2.2)
    width = Inches(4)
    height = Inches(1.8)
    
    feature_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = feature_box.text_frame
    
    # Feature title
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Feature description
    p = text_frame.add_paragraph()
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(8)

# Slide 5: Key Innovation
items = [
    "✓ Pattern Recognition: Advanced regex for character substitution detection",
    "✓ VirusTotal Integration: Real-time threat intelligence from 70+ scanners",
    "✓ AI Analysis: Gemini AI for contextual threat assessment",
    "✓ Multi-Modal Scanning: QR codes, screenshots, and OCR analysis"
]
add_content_slide(prs, "🔍 Key Innovation\nMulti-Modal AI Phishing Detection", items, (250, 112, 154))

# Slide 6: Technical Architecture
items = [
    "✓ Frontend: Next.js 14 with App Router for optimal performance",
    "✓ AI Engine: Google Gemini for natural language processing",
    "✓ Security APIs: VirusTotal for real-time threat intelligence",
    "✓ Image Processing: Tesseract.js for OCR, jsQR for QR code scanning",
    "✓ Deployment: Vercel for serverless scalability"
]
add_content_slide(prs, "🏗️ Technical Architecture", items, (48, 207, 208))

# Slide 7: Implementation Approach
items = [
    "✓ Phase 1: Core infrastructure and design system (30 min)",
    "✓ Phase 2: Phishing detection engine - Full implementation (90 min)",
    "✓ Phase 3: Awareness, incident response, code security (2 hours)",
    "✓ Phase 4: Testing and refinement (15 min)",
    "✓ Phase 5: Deployment to Vercel (10 min)",
    "",
    "Total Build Time: ~6 hours | Result: Demo-ready MVP"
]
add_content_slide(prs, "⚙️ Implementation Strategy", items, (168, 237, 234))

# Slide 8: Features Deep Dive
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 154, 86)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "✨ Feature Highlights"
p = title_frame.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Features
features = [
    ("🎮 Gamified Learning", "Interactive quizzes with instant feedback"),
    ("⚡ Real-time Analysis", "Instant URL scanning with risk scoring"),
    ("📋 Smart Playbooks", "Guided response for multiple threats"),
    ("🔐 Code Scanning", "Detect SQL injection, XSS, secrets")
]

for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 4.7)
    top = Inches(2.3 + row * 2.2)
    width = Inches(4)
    height = Inches(1.8)
    
    feature_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = feature_box.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = text_frame.add_paragraph()
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(8)

# Slide 9: Impact & Scalability
items = [
    "✓ Accessibility: Free tier for individuals and small businesses",
    "✓ Scalability: Serverless architecture handles millions of users",
    "✓ Education: Makes cybersecurity knowledge accessible to all",
    "✓ Prevention: Reduces phishing success rate by early detection",
    "✓ Cost Savings: Prevents costly security incidents"
]
add_content_slide(prs, "🚀 Impact & Scalability", items, (106, 17, 203))

# Slide 10: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(11, 163, 96)

# Emoji
emoji_box = slide.shapes.add_textbox(Inches(4), Inches(1.5), Inches(2), Inches(1))
emoji_frame = emoji_box.text_frame
emoji_frame.text = "🎉"
p = emoji_frame.paragraphs[0]
p.font.size = Pt(80)
p.alignment = PP_ALIGN.CENTER

# Thank you
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Thank You!"
p = title_frame.paragraphs[0]
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Making Cybersecurity Accessible to Everyone"
p = subtitle_frame.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Closing
closing_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(7), Inches(1.5))
closing_frame = closing_box.text_frame
closing_frame.text = "CyberGuardian AI\nLive Demo Available\n\nQuestions? Let's discuss how we can make the digital world safer."
p = closing_frame.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

for paragraph in closing_frame.paragraphs[1:]:
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    paragraph.alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('CyberGuardian_Presentation.pptx')
print("PowerPoint presentation created successfully!")
print("File: CyberGuardian_Presentation.pptx")
